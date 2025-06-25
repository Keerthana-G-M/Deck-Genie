import io
import streamlit as st
import hashlib
from typing import Dict, Any, List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from image_fetcher import fetch_image_for_slide, get_slide_icon
from utils import FONTS, COLORS, FONT_SIZES, MARGINS, CONTENT_AREA, IMAGE_AREA, match_icon_to_feature

# Enhanced presentation styles with persona-specific configurations
PRESENTATION_STYLES = {
    "technical": {
        "colors": {
            "primary": "2F2F2F",    # Dark gray
            "secondary": "0096C7",   # Bright blue
            "accent": "FF6B6B",     # Coral
            "text": "2F2F2F",       # Dark gray
            "background": "FFFFFF",  # White
            "bullet": "0096C7"      # Bullet point color
        },
        "fonts": {
            "heading": "Calibri",
            "body": "Calibri"
        },
        "font_sizes": {
            "title": 32,      # Main slide title
            "subtitle": 24,   # Subtitle
            "heading": 20,    # Section headers
            "body": 14,       # Regular content
            "bullet": 14,     # Bullet points
            "caption": 12     # Small text/captions
        },
        "spacing": {
            "line_spacing": 1.15,
            "paragraph_spacing": 6,
            "bullet_indent": 0.25
        }
    },
    "executive": {
        "colors": {
            "primary": "1F497D",    # Deep blue
            "secondary": "4F81BD",   # Light blue
            "accent": "C0504D",     # Burgundy
            "text": "000000",       # Black
            "background": "FFFFFF",  # White
            "bullet": "4472C4"      # Bullet point color
        },
        "fonts": {
            "heading": "Calibri",
            "body": "Calibri"
        },
        "font_sizes": {
            "title": 36,      # Larger title for impact
            "subtitle": 28,   # Larger subtitle
            "heading": 24,    # Prominent headers
            "body": 16,       # Larger body text
            "bullet": 16,     # Larger bullets
            "caption": 14     # Larger captions
        },
        "spacing": {
            "line_spacing": 1.2,
            "paragraph_spacing": 8,
            "bullet_indent": 0.3
        }
    },
    "business": {
        "colors": {
            "primary": "2F5597",    # Warm blue
            "secondary": "70AD47",   # Green
            "accent": "FFC000",     # Gold
            "text": "2F2F2F",       # Dark gray
            "background": "FFFFFF",  # White
            "bullet": "70AD47"      # Bullet point color
        },
        "fonts": {
            "heading": "Calibri",
            "body": "Calibri"
        },
        "font_sizes": {
            "title": 32,      # Standard title
            "subtitle": 26,   # Clear subtitle
            "heading": 22,    # Clear headers
            "body": 15,       # Readable body
            "bullet": 15,     # Matching bullets
            "caption": 13     # Readable captions
        },
        "spacing": {
            "line_spacing": 1.15,
            "paragraph_spacing": 7,
            "bullet_indent": 0.25
        }
    }
}

# Persona-specific content templates
PERSONA_CONTENT = {
    "technical": {
        "problem_slide": {
            "title": "Technical Challenges",
            "bullets": [
                "Complex integration requirements with legacy systems",
                "Scalability limitations in current architecture",
                "Manual deployment and configuration processes",
                "Limited visibility into system performance metrics",
                "Technical debt from outdated infrastructure",
                "Resource-intensive maintenance procedures"
            ]
        },
        "solution_slide": {
            "title": "Technical Solution Architecture",
            "description": "Our solution leverages modern microservices architecture with containerized deployments for optimal scalability and maintainability.",
            "features": [
                "Containerized microservices deployment",
                "RESTful API integration capabilities",
                "Automated CI/CD pipeline integration",
                "Real-time monitoring and logging",
                "Horizontal scaling with Kubernetes"
            ]
        },
        "advantage_slide": {
            "title": "Technical Advantages",
            "bullets": [
                "99.99% system availability with redundant architecture",
                "Sub-millisecond response times for API endpoints",
                "Zero-downtime deployments with blue-green strategy",
                "Comprehensive API documentation and SDKs",
                "Advanced debugging and monitoring tools"
            ]
        },
        "cta_slide": {
            "title": "Ready to Enhance Your Technical Infrastructure?",
            "cta_text": "Schedule a technical deep-dive session today",
            "contact_info": "Contact our solutions architects to begin"
        }
    },
    "executive": {
        "problem_slide": {
            "title": "Strategic Business Challenges",
            "bullets": [
                "Increasing operational costs impacting bottom line",
                "Market share erosion from digital competitors",
                "Risk exposure from legacy systems",
                "Inefficient resource allocation",
                "Compliance and regulatory pressures",
                "Limited business agility and innovation"
            ]
        },
        "solution_slide": {
            "title": "Strategic Business Solution",
            "description": "Our enterprise solution delivers measurable ROI through operational efficiency, risk reduction, and enhanced business capabilities.",
            "features": [
                "Comprehensive business intelligence dashboard",
                "Executive-level reporting and analytics",
                "Strategic resource optimization",
                "Risk management framework",
                "Compliance automation suite"
            ]
        },
        "advantage_slide": {
            "title": "Business Impact & ROI",
            "bullets": [
                "30% reduction in operational costs",
                "40% improvement in resource utilization",
                "60% faster time-to-market for new initiatives",
                "25% increase in customer satisfaction",
                "Demonstrated ROI within 6 months"
            ]
        },
        "cta_slide": {
            "title": "Ready to Drive Business Growth?",
            "cta_text": "Transform your business today",
            "contact_info": "Contact our executive team to discuss your goals"
        }
    },
    "business": {
        "problem_slide": {
            "title": "Business Operations Challenges",
            "bullets": [
                "Inefficient workflow processes causing delays",
                "Data silos limiting business insights",
                "Manual reporting consuming valuable time",
                "Inconsistent customer experience",
                "Resource allocation inefficiencies",
                "Limited scalability for growth"
            ]
        },
        "solution_slide": {
            "title": "Business Process Solution",
            "description": "Our solution streamlines operations, automates workflows, and provides actionable insights for better business decisions.",
            "features": [
                "Automated workflow management",
                "Intuitive business process designer",
                "Real-time performance analytics",
                "Customizable reporting dashboard",
                "Integration with existing tools"
            ]
        },
        "advantage_slide": {
            "title": "Operational Benefits",
            "bullets": [
                "Streamlined business processes",
                "Enhanced operational visibility",
                "Improved team collaboration",
                "Reduced manual workload",
                "Better resource utilization"
            ]
        },
        "cta_slide": {
            "title": "Ready to Streamline Your Operations?",
            "cta_text": "Start optimizing your business processes",
            "contact_info": "Contact us to begin your transformation"
        }
    }
}

# Standard slide dimensions and margins for all personas
SLIDE_MARGINS = {
    'left': Inches(0.5),
    'right': Inches(0.5),
    'top': Inches(0.2),
    'bottom': Inches(0.5)
}

# Content dimensions adjusted per persona
def get_content_dimensions(style_config):
    """Get content dimensions based on persona style."""
    return {
        'main_content': {
            'height': Inches(1.2),
            'font_size': Pt(style_config["font_sizes"]["body"])
        },
        'bullet_section': {
            'height': Inches(3.8),
            'font_size': Pt(style_config["font_sizes"]["bullet"]),
            'header_size': Pt(style_config["font_sizes"]["heading"])
        },
        'image_section': {
            'width': Inches(4.5),
            'height': Inches(4.0)
        }
    }

def apply_text_formatting(text_frame, style_config, is_title=False, is_bullet=False, is_header=False):
    """Apply consistent text formatting based on persona style."""
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT if is_bullet else PP_ALIGN.CENTER
        paragraph.line_spacing = style_config["spacing"]["line_spacing"]
        
        if is_bullet:
            paragraph.space_before = Pt(style_config["spacing"]["paragraph_spacing"])
            paragraph.space_after = Pt(style_config["spacing"]["paragraph_spacing"])
            paragraph.left_indent = Inches(style_config["spacing"]["bullet_indent"])
        
        for run in paragraph.runs:
            run.font.name = style_config["fonts"]["heading"] if (is_title or is_header) else style_config["fonts"]["body"]
            
            if is_title:
                run.font.size = Pt(style_config["font_sizes"]["title"])
            elif is_header:
                run.font.size = Pt(style_config["font_sizes"]["heading"])
            elif is_bullet:
                run.font.size = Pt(style_config["font_sizes"]["bullet"])
            else:
                run.font.size = Pt(style_config["font_sizes"]["body"])
            
            run.font.color.rgb = RGBColor.from_string(
                style_config["colors"]["primary"] if (is_title or is_header) else style_config["colors"]["text"]
            )

def create_section_header(slide, title, style="business"):
    """Create an enhanced section header with improved styling and spacing."""
    style = normalize_style(style)
    style_config = PRESENTATION_STYLES[style]
    
    # Add title with adjusted positioning
    title_left = Inches(0.5)
    title_top = Inches(0.2)
    title_width = Inches(12.33)
    title_height = Inches(0.6)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    title_p = title_tf.paragraphs[0]
    
    # Clean up title text
    clean_title = " ".join(title.split())
    title_p.text = clean_title
    title_p.alignment = PP_ALIGN.LEFT
    
    # Apply enhanced title formatting with style-specific font size
    for run in title_p.runs:
        run.font.name = style_config["fonts"]["heading"]
        run.font.size = Pt(style_config["font_sizes"]["heading"])
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(style_config["colors"]["primary"])
    
    # Add accent line with enhanced styling
    line_left = Inches(0.5)
    line_top = title_top + Inches(0.5)
    line_width = Inches(12.33)
    line_height = Inches(0.02)
    
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, line_left, line_top, line_width, line_height)
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor.from_string(style_config["colors"]["accent"])
    line.line.fill.background()
    
    return line_top + Inches(0.1)

def deduplicate_content(content_list, product_name=None):
    """Helper function to deduplicate content while preserving formatting."""
    seen_points = {}  # Dictionary to track unique content
    cleaned_content = []
    
    for item in content_list:
        # Handle both dictionary and string items
        if isinstance(item, dict):
            text = item.get("point", "") or item.get("feature", "") or item.get("text", "") or str(item)
        else:
            text = str(item)
        
        # Clean the text
        clean_text = clean_and_deduplicate_text(text)
        if not clean_text:  # Skip empty content
            continue
            
        # Get comparison key
        key = get_comparison_key(clean_text)
        if not key:  # Skip if key is empty
            continue
            
        # If this is a new unique point or longer version of existing point
        if key not in seen_points or len(clean_text) > len(seen_points[key]):
            seen_points[key] = clean_text
            if isinstance(item, dict):
                # Preserve the dictionary structure but update the text
                item_copy = item.copy()
                for k, v in item_copy.items():
                    if isinstance(v, str):
                        item_copy[k] = v.replace("[Product Name]", product_name) if product_name else v
                cleaned_content.append(item_copy)
            else:
                # Replace product name if present
                clean_point = clean_text.replace("[Product Name]", product_name) if product_name else clean_text
                cleaned_content.append(clean_point)
    
    return cleaned_content

def clean_and_deduplicate_text(text):
    """Clean and prepare text for deduplication."""
    if not text:
        return ""
        
    # Remove common bullet points and markers
    text = text.replace('•', '').replace('*', '').replace('·', '').replace('-', '')
    
    # Remove extra whitespace
    text = ' '.join(text.split())
    
    # Remove trailing punctuation
    text = text.rstrip('.,;:')
    
    return text.strip()

def get_comparison_key(text):
    """Create a standardized key for text comparison."""
    if not text:
        return ""
        
    # Convert to lowercase and remove punctuation
    key = ''.join(c.lower() for c in text if c.isalnum() or c.isspace())
    
    # Remove common words that shouldn't affect uniqueness
    stop_words = {'the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'with', 'by'}
    words = key.split()
    key_words = [w for w in words if w not in stop_words]
    
    return ' '.join(key_words)

def deduplicate_bullets(bullets):
    """Remove duplicate bullet points while preserving the best version of each."""
    if not bullets:
        return []
        
    unique_bullets = {}
    
    for bullet in bullets:
        # Handle both string and dictionary bullets
        if isinstance(bullet, dict):
            text = bullet.get('point', '') or bullet.get('text', '') or str(bullet)
        else:
            text = str(bullet)
            
        # Clean and get comparison key
        clean_text = clean_and_deduplicate_text(text)
        key = get_comparison_key(clean_text)
        
        # Skip empty content
        if not key:
            continue
            
        # Keep the longer version of duplicate content
        if key not in unique_bullets or len(clean_text) > len(unique_bullets[key]):
            unique_bullets[key] = clean_text
    
    return list(unique_bullets.values())

def create_problem_slide(prs: Presentation, content: Dict[str, Any], presentation_context: Dict[str, Any] = None, used_images: set = None):
    """Create the problem statement slide with enhanced styling and proper spacing."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    style = presentation_context.get('style', 'business')
    style_config = PRESENTATION_STYLES[style]
    content_dimensions = get_content_dimensions(style_config)
    product_name = presentation_context.get("metadata", {}).get("product_name", "")
    
    # Create section header with accent line and get the content start position
    clean_title = content["title"].replace("[Product Name]", product_name)
    # Remove any HTML tags from title
    clean_title = clean_title.replace("<strong>", "").replace("</strong>", "")
    content_start_y = create_section_header(slide, clean_title, style)
    
    # Collect all bullet points
    all_points = []
    
    # Extract points from all possible sources and clean HTML tags
    if "pain_points" in content:
        all_points.extend([p.replace("<strong>", "").replace("</strong>", "") for p in content["pain_points"]])
    if "bullets" in content:
        all_points.extend([b.replace("<strong>", "").replace("</strong>", "") for b in content["bullets"]])
    if "differentiators" in content:
        for item in content["differentiators"]:
            if isinstance(item, dict):
                point = item.get("point", "") or item.get("text", "") or str(item)
            else:
                point = str(item)
            if point:
                all_points.append(point.replace("<strong>", "").replace("</strong>", ""))
    
    # Deduplicate and clean all points
    clean_bullets = deduplicate_bullets(all_points)
    
    # Replace product name placeholders
    if product_name:
        clean_bullets = [bullet.replace("[Product Name]", product_name) for bullet in clean_bullets]
    
    # Ensure minimum content
    while len(clean_bullets) < 4:
        new_bullet = generate_additional_bullet(determine_content_type(clean_bullets))
        if get_comparison_key(new_bullet) not in {get_comparison_key(b) for b in clean_bullets}:
            clean_bullets.append(new_bullet)
    
    # Limit to maximum bullets and truncate each bullet based on persona
    max_bullets = 4 if style == "executive" else 5
    clean_bullets = clean_bullets[:max_bullets]
    
    # Adjust truncation length based on persona
    max_chars = 100 if style == "executive" else 80
    max_words = 20 if style == "executive" else 15
    clean_bullets = [truncate_text_for_slide(bullet, max_chars=max_chars, max_words=max_words) 
                    for bullet in clean_bullets]
    
    # Create text box for bullets with adjusted dimensions
    text_box = slide.shapes.add_textbox(
        SLIDE_MARGINS['left'],
        content_start_y + Inches(0.3),
        Inches(6.8),
        content_dimensions['bullet_section']['height']
    )
    
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    
    # Add bullets with persona-specific formatting
    for i, bullet_text in enumerate(clean_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = f"• {bullet_text}"
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(style_config["spacing"]["paragraph_spacing"])
        p.space_after = Pt(style_config["spacing"]["paragraph_spacing"])
        p.line_spacing = style_config["spacing"]["line_spacing"]
        
        # Apply consistent formatting based on persona
        for run in p.runs:
            run.font.name = style_config["fonts"]["body"]
            run.font.size = Pt(style_config["font_sizes"]["bullet"])
            run.font.color.rgb = RGBColor.from_string(style_config["colors"]["text"])
    
    # Add image with adjusted positioning
    if used_images is None:
        used_images = set()
    
    image_data = fetch_image_with_cache("problem", presentation_context, used_images, use_placeholders=False)
    add_slide_image(slide, image_data, content_start_y, style_config)
    
    return slide

def normalize_style(style: str) -> str:
    """Normalize style/persona name to lowercase and validate it exists."""
    if not isinstance(style, str):
        return 'business'
    
    style = style.lower()
    return style if style in PRESENTATION_STYLES else 'business'

def create_title_slide(prs: Presentation, content: Dict[str, str], style="business"):
    """Create an enhanced title slide with styling."""
    style = normalize_style(style)
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    style_config = PRESENTATION_STYLES[style]
    
    # Add large title with adjusted positioning
    title_left = Inches(1.0)
    title_top = Inches(2.0)  # Moved up from 2.5
    title_width = Inches(11.33)
    title_height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_p = title_tf.paragraphs[0]
    title_p.text = content.get("title", "Presentation Title")
    title_p.alignment = PP_ALIGN.CENTER
    
    for run in title_p.runs:
        run.font.name = style_config["fonts"]["heading"]
        run.font.size = Pt(44)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(style_config["colors"]["primary"])
    
    # Add subtitle if present with adjusted spacing
    if "subtitle" in content:
        subtitle_top = title_top + Inches(2.0)  # Increased spacing from 1.8
        subtitle_height = Inches(0.8)
        
        subtitle_box = slide.shapes.add_textbox(title_left, subtitle_top, title_width, subtitle_height)
        subtitle_tf = subtitle_box.text_frame
        subtitle_tf.word_wrap = True
        subtitle_p = subtitle_tf.paragraphs[0]
        subtitle_p.text = content["subtitle"]
        subtitle_p.alignment = PP_ALIGN.CENTER
        
        for run in subtitle_p.runs:
            run.font.name = style_config["fonts"]["body"]
            run.font.size = Pt(28)
            run.font.color.rgb = RGBColor.from_string(style_config["colors"]["secondary"])
    
    # Add decorative accent line with adjusted position
    line_left = Inches(2.0)
    line_top = subtitle_top + Inches(1.0)  # Position relative to subtitle with more space
    line_width = Inches(9.33)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, line_left, line_top, line_width, Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor.from_string(style_config["colors"]["accent"])
    line.line.fill.background()

def create_solution_slide(prs: Presentation, content: Dict[str, Any], presentation_context: Dict[str, Any] = None, used_images: set = None):
    """Create the solution overview slide with enhanced content."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    style = presentation_context.get('style', 'business')
    style_config = PRESENTATION_STYLES[style]
    content_dimensions = get_content_dimensions(style_config)
    product_name = presentation_context.get("metadata", {}).get("product_name", "")
    
    # Create section header
    clean_title = content["title"].replace("[Product Name]", product_name) if product_name else content["title"]
    # Remove any formatting markers that might be in the title
    clean_title = clean_title.replace("{heading{", "").replace("}}", "")
    content_start_y = create_section_header(slide, clean_title, style)
    
    # Get content text and clean it
    content_text = content.get("paragraph", content.get("description", content.get("value_proposition", "")))
    clean_text = content_text.replace("[Product Name]", product_name) if product_name else content_text
    # Remove any formatting markers from the content
    clean_text = clean_text.replace("{heading{", "").replace("}}", "")
    
    # Create main content section with standardized dimensions
    main_box = slide.shapes.add_textbox(
        SLIDE_MARGINS['left'],
        content_start_y + SLIDE_MARGINS['top'],
        Inches(6.8),
        content_dimensions['main_content']['height']
    )
    
    main_tf = main_box.text_frame
    main_tf.word_wrap = True
    main_tf.margin_left = Inches(0.1)
    main_tf.margin_right = Inches(0.1)
    
    # Add main description with standardized formatting
    main_p = main_tf.paragraphs[0]
    main_p.text = truncate_text_for_slide(clean_text, max_chars=250, max_words=50)
    main_p.alignment = PP_ALIGN.LEFT
    
    for run in main_p.runs:
        run.font.name = style_config["fonts"]["body"]
        run.font.size = Pt(style_config["font_sizes"]["body"])
        run.font.color.rgb = RGBColor.from_string(style_config["colors"]["text"])
    
    # Create features section with standardized dimensions
    features_box = slide.shapes.add_textbox(
        SLIDE_MARGINS['left'],
        content_start_y + Inches(1.8),
        Inches(6.8),
        content_dimensions['bullet_section']['height']
    )
    
    features_tf = features_box.text_frame
    features_tf.word_wrap = True
    
    # Add features header with clean text
    features_header = features_tf.paragraphs[0]
    features_header.text = "Key Solution Features"  # Simplified header text
    features_header.alignment = PP_ALIGN.LEFT
    features_header.space_after = Pt(style_config["spacing"]["paragraph_spacing"])
    
    for run in features_header.runs:
        run.font.name = style_config["fonts"]["heading"]
        run.font.size = Pt(style_config["font_sizes"]["heading"])
        run.font.color.rgb = RGBColor.from_string(style_config["colors"]["primary"])
    
    # Add feature points with standardized formatting
    feature_points = [
        "Advanced threat detection with real-time monitoring",
        "Automated incident response and remediation",
        "Seamless integration with existing security tools",
        "Comprehensive analytics and reporting dashboard",
        "24/7 expert support and continuous updates"
    ]
    
    # Adjust number of features based on persona
    max_features = 4 if style == "executive" else 5
    feature_points = feature_points[:max_features]
    
    for point in feature_points:
        p = features_tf.add_paragraph()
        # Adjust truncation based on persona
        max_chars = 100 if style == "executive" else 80
        max_words = 20 if style == "executive" else 15
        truncated_point = truncate_text_for_slide(point, max_chars=max_chars, max_words=max_words)
        p.text = f"• {truncated_point}"
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(style_config["spacing"]["paragraph_spacing"])
        p.space_after = Pt(style_config["spacing"]["paragraph_spacing"])
        p.line_spacing = style_config["spacing"]["line_spacing"]
        
        for run in p.runs:
            run.font.name = style_config["fonts"]["body"]
            run.font.size = Pt(style_config["font_sizes"]["bullet"])
            run.font.color.rgb = RGBColor.from_string(style_config["colors"]["text"])
    
    # Add image with standardized dimensions
    if used_images is None:
        used_images = set()
    
    image_data = fetch_image_with_cache("solution", presentation_context, used_images, use_placeholders=False)
    add_slide_image(slide, image_data, content_start_y, style_config)
    
    return slide

def create_features_slide(prs: Presentation, content: Dict[str, Any], style="business"):
    """Create the key features slide with enhanced styling and deduplication."""
    style = normalize_style(style)
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    style_config = PRESENTATION_STYLES[style]
    
    # Create section header with accent line
    clean_title = content["title"].replace("[Product Name]", "")
    create_section_header(slide, clean_title, style)
    
    # Get features list and deduplicate
    features = deduplicate_content(content.get("features", []))
    max_features = 5
    features = features[:max_features]
    
    # Calculate spacing based on number of features
    if len(features) <= 3:
        feature_height = Inches(1.5)
    else:
        feature_height = Inches(1.2)
    
    # Create feature list with icons
    left_margin = Inches(1.0)
    top_start = Inches(1.7)
    icon_width = Inches(0.8)
    text_width = Inches(11.0)
    
    for i, feature in enumerate(features):
        # Get feature text
        if isinstance(feature, dict):
            feature_text = feature.get('feature', feature.get('name', feature.get('title', str(feature))))
        else:
            feature_text = str(feature)
        
        # Clean and prepare text
        feature_text = truncate_text_for_slide(feature_text, max_chars=100, max_words=20)
        
        # Calculate position
        top_position = top_start + (i * feature_height)
        
        # Add icon with styling
        icon_box = slide.shapes.add_textbox(left_margin, top_position, icon_width, feature_height)
        icon_tf = icon_box.text_frame
        icon_tf.auto_size = False
        icon_p = icon_tf.paragraphs[0]
        icon_p.text = match_icon_to_feature(feature_text)
        icon_p.alignment = PP_ALIGN.CENTER
        
        # Apply icon formatting
        for run in icon_p.runs:
            run.font.name = style_config["fonts"]["body"]
            run.font.size = Pt(24)
            run.font.bold = True
            run.font.color.rgb = RGBColor.from_string(style_config["colors"]["primary"])
        
        # Add feature text with styling
        text_left = left_margin + icon_width + Inches(0.3)
        text_box = slide.shapes.add_textbox(text_left, top_position, text_width, feature_height)
        text_tf = text_box.text_frame
        text_tf.word_wrap = True
        text_tf.auto_size = False
        text_p = text_tf.paragraphs[0]
        text_p.text = feature_text
        text_p.alignment = PP_ALIGN.LEFT
        
        # Apply text formatting
        text_p.space_before = Pt(3)
        text_p.space_after = Pt(3)
        text_p.line_spacing = style_config["spacing"]["line_spacing"]
        
        for run in text_p.runs:
            run.font.name = style_config["fonts"]["body"]
            run.font.size = Pt(20)
            run.font.color.rgb = RGBColor.from_string(style_config["colors"]["text"])

def expand_advantage_content(text, product_name=""):
    """Expand competitive advantage content with specific details."""
    advantage_expansions = [
        "This provides a clear competitive edge in the market.",
        "Our solution outperforms traditional approaches significantly.",
        "Customers report substantial improvements in key metrics.",
        "Implementation time and complexity are reduced by up to 60%.",
        "The return on investment is typically realized within 6 months.",
        "Our approach has been validated by industry experts and analysts.",
        "Integration capabilities exceed industry standards.",
        "Customer satisfaction scores consistently exceed 90%."
    ]
    
    clean_text = text.strip()
    if product_name:
        clean_text = clean_text.replace("[Product Name]", product_name)
    
    # Add expansions if content is too short
    if len(clean_text) < 200:
        for expansion in advantage_expansions:
            if len(clean_text) >= 200:
                break
            if not clean_text.endswith('.'):
                clean_text += "."
            clean_text += " " + expansion
    
    return clean_text

def expand_audience_content(text, product_name=""):
    """Expand target audience content with detailed insights."""
    audience_expansions = [
        "These organizations face increasing pressure to modernize their operations.",
        "Decision-makers in this segment prioritize efficiency and reliability.",
        "The need for comprehensive solutions continues to grow.",
        "Budget constraints and resource limitations drive demand for efficient solutions.",
        "Regulatory compliance requirements add complexity to their operations.",
        "They seek partners who understand their unique challenges.",
        "Time-to-value is a critical factor in their decision-making process.",
        "Integration with existing systems is a key consideration."
    ]
    
    clean_text = text.strip()
    if product_name:
        clean_text = clean_text.replace("[Product Name]", product_name)
    
    # Add expansions if content is too short
    if len(clean_text) < 200:
        for expansion in audience_expansions:
            if len(clean_text) >= 200:
                break
            if not clean_text.endswith('.'):
                clean_text += "."
            clean_text += " " + expansion
    
    return clean_text

def expand_market_content(text):
    """Expand market analysis content with detailed insights."""
    market_expansions = [
        "Market growth is driven by increasing digital transformation initiatives.",
        "Industry analysts project significant expansion in coming years.",
        "Regulatory requirements continue to drive market demand.",
        "Technology adoption rates show strong upward trends.",
        "Organizations are increasing their investment in this sector.",
        "Market penetration opportunities remain substantial.",
        "Competitive dynamics favor innovative solutions.",
        "Global market conditions support continued growth."
    ]
    
    clean_text = text.strip()
    
    # Add expansions if content is too short
    if len(clean_text) < 200:
        for expansion in market_expansions:
            if len(clean_text) >= 200:
                break
            if not clean_text.endswith('.'):
                clean_text += "."
            clean_text += " " + expansion
    
    return clean_text

def create_advantage_slide(prs: Presentation, content: Dict[str, Any], presentation_context: Dict[str, Any] = None, used_images: set = None):
    """Create the competitive advantage slide with enhanced content."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    product_name = presentation_context.get("metadata", {}).get("product_name", "")
    style = presentation_context.get('style', 'formal')
    style_config = PRESENTATION_STYLES[style]
    
    # Create section header
    clean_title = content["title"].replace("[Product Name]", product_name)
    content_start_y = create_section_header(slide, clean_title, style)
    
    # Get and enhance bullet points
    bullets = []
    if "differentiators" in content:
        bullets = content["differentiators"]
    elif "bullets" in content:
        bullets = content["bullets"]
    elif "advantages" in content:
        bullets = content["advantages"]
        
    # Add only the most important additional advantages
    additional_advantages = [
        "Industry-leading innovation and technology",
        "Proven track record of success",
        "Comprehensive support and resources"
    ]
    bullets.extend(additional_advantages)
    
    # Create bullet list with focused content
    bullet_box = slide.shapes.add_textbox(
        Inches(0.5),
        content_start_y + Inches(0.2),
        Inches(6.8),
        Inches(4.5)
    )
    
    tf = bullet_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    
    # Clean and limit bullets
    clean_bullets = deduplicate_bullets(bullets)
    clean_bullets = clean_bullets[:5]  # Limit to 5 key points
    
    # Add bullets with consistent formatting
    for i, bullet_text in enumerate(clean_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            
        # Keep bullet text concise
        truncated_text = truncate_text_for_slide(bullet_text, max_chars=80, max_words=15)
        p.text = f"• {truncated_text}"
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(12)
        p.space_after = Pt(12)
        
        # Apply consistent formatting (no bold)
        for run in p.runs:
            run.font.name = style_config["fonts"]["body"]
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor.from_string(style_config["colors"]["text"])
    
    # Add image with proper positioning
    if used_images is None:
        used_images = set()
    
    image_data = fetch_image_with_cache("advantage", presentation_context, used_images, use_placeholders=False)
    
    if image_data:
        pic = slide.shapes.add_picture(
            image_data,
            Inches(7.8),
            content_start_y + Inches(0.2),
            Inches(4.5),
            Inches(4.0)
        )
    else:
        # Add styled fallback icon
        icon_left = Inches(9.5)
        icon_top = content_start_y + Inches(1.0)
        icon_box = slide.shapes.add_textbox(icon_left, icon_top, Inches(2.0), Inches(2.0))
        icon_tf = icon_box.text_frame
        icon_tf.auto_size = False
        icon_p = icon_tf.paragraphs[0]
        icon_p.text = get_slide_icon("advantage")["icon"]
        icon_p.alignment = PP_ALIGN.CENTER
        
        for run in icon_p.runs:
            run.font.name = style_config["fonts"]["body"]
            run.font.size = Pt(64)
            run.font.color.rgb = RGBColor.from_string(style_config["colors"]["primary"])
            
    # Add concise bottom message
    bottom_box = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(6.5),
        Inches(12.33),
        Inches(0.5)
    )
    bottom_tf = bottom_box.text_frame
    bottom_p = bottom_tf.paragraphs[0]
    bottom_p.text = "Partner with us to transform your business today."
    bottom_p.alignment = PP_ALIGN.CENTER
    
    for run in bottom_p.runs:
        run.font.name = style_config["fonts"]["body"]
        run.font.size = Pt(16)
        run.font.italic = True
        run.font.color.rgb = RGBColor.from_string(style_config["colors"]["accent"])

def create_audience_slide(prs: Presentation, content: Dict[str, Any], presentation_context: Dict[str, Any] = None, used_images: set = None):
    """Create the target audience slide with enhanced content."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    product_name = presentation_context.get("metadata", {}).get("product_name", "")
    style = presentation_context.get('style', 'formal')
    style_config = PRESENTATION_STYLES[style]
    
    # Create section header
    clean_title = content["title"].replace("[Product Name]", product_name)
    content_start_y = create_section_header(slide, clean_title, style)
    
    # Get and enhance content
    content_text = content.get("paragraph", content.get("description", content.get("content", "")))
    
    # Add audience-specific insights
    audience_insights = [
        "Key decision-makers in enterprise security",
        "IT and security operations teams",
        "Compliance and risk management professionals",
        "Technology leaders driving digital transformation"
    ]
    
    # Create main content section
    main_box = slide.shapes.add_textbox(
        Inches(0.5),
        content_start_y + Inches(0.2),
        Inches(6.8),
        Inches(2.2)  # Reduced from 2.5 to prevent overlap
    )
    
    main_tf = main_box.text_frame
    main_tf.word_wrap = True
    main_tf.margin_left = Inches(0.1)
    main_tf.margin_right = Inches(0.1)
    
    # Add enhanced main content
    expanded_content = expand_audience_content(content_text, product_name)
    main_p = main_tf.paragraphs[0]
    main_p.text = expanded_content
    main_p.alignment = PP_ALIGN.LEFT
    
    for run in main_p.runs:
        run.font.name = style_config["fonts"]["body"]
        run.font.size = Pt(14)  # Reduced from 16 to fit better
        run.font.color.rgb = RGBColor.from_string(style_config["colors"]["text"])
    
    # Add audience segments section with adjusted position
    segments_box = slide.shapes.add_textbox(
        Inches(0.5),
        content_start_y + Inches(2.6),  # Moved down from 3.0
        Inches(6.8),
        Inches(2.5)
    )
    
    segments_tf = segments_box.text_frame
    segments_tf.word_wrap = True
    
    # Add header for segments
    header_p = segments_tf.paragraphs[0]
    header_p.text = "Key Audience Segments:"
    header_p.alignment = PP_ALIGN.LEFT
    
    for run in header_p.runs:
        run.font.name = style_config["fonts"]["body"]
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(style_config["colors"]["primary"])
    
    # Add audience segments with enhanced formatting
    for insight in audience_insights:
        p = segments_tf.add_paragraph()
        p.text = f"• {insight}"
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        
        for run in p.runs:
            run.font.name = style_config["fonts"]["body"]
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor.from_string(style_config["colors"]["text"])
    
    # Add image with proper positioning
    if used_images is None:
        used_images = set()
    
    image_data = fetch_image_with_cache("audience", presentation_context, used_images, use_placeholders=False)
    
    if image_data:
        pic = slide.shapes.add_picture(
            image_data,
            Inches(7.8),
            content_start_y + Inches(0.2),
            Inches(4.5),
            Inches(4.5)
        )
    else:
        # Add styled fallback icon
        icon_left = Inches(9.5)
        icon_top = content_start_y + Inches(1.0)
        icon_box = slide.shapes.add_textbox(icon_left, icon_top, Inches(2.0), Inches(2.0))
        icon_tf = icon_box.text_frame
        icon_tf.auto_size = False
        icon_p = icon_tf.paragraphs[0]
        icon_p.text = get_slide_icon("audience")["icon"]
        icon_p.alignment = PP_ALIGN.CENTER
        
        for run in icon_p.runs:
            run.font.name = style_config["fonts"]["body"]
            run.font.size = Pt(64)
            run.font.color.rgb = RGBColor.from_string(style_config["colors"]["primary"])
    
    # Add value proposition message with adjusted position
    value_box = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(6.8),  # Moved down from 6.5
        Inches(12.33),
        Inches(0.5)
    )
    value_tf = value_box.text_frame
    value_p = value_tf.paragraphs[0]
    value_p.text = "We understand your unique challenges and are ready to help you succeed."
    value_p.alignment = PP_ALIGN.CENTER
    
    for run in value_p.runs:
        run.font.name = style_config["fonts"]["body"]
        run.font.size = Pt(16)
        run.font.italic = True
        run.font.color.rgb = RGBColor.from_string(style_config["colors"]["accent"])

def expand_cta_content(text, content_type="cta"):
    """Expand call-to-action content with compelling details."""
    cta_expansions = [
        "Take the first step toward transforming your security posture.",
        "Join industry leaders who have already embraced this solution.",
        "Schedule a personalized demo to see the benefits firsthand.",
        "Our team of experts is ready to guide your implementation.",
        "Start your journey to enhanced security today.",
        "Limited-time implementation support available.",
        "Join our growing community of satisfied customers.",
        "Experience the difference our solution can make."
    ]
    
    clean_text = text.strip()
    
    # Add expansions if content is too short
    if len(clean_text) < 100:  # CTAs can be shorter
        for expansion in cta_expansions:
            if len(clean_text) >= 100:
                break
            if not clean_text.endswith('.'):
                clean_text += "."
            clean_text += " " + expansion
    
    return clean_text

def create_cta_slide(prs: Presentation, content: Dict[str, Any], style="business"):
    """Create a persona-specific call to action slide."""
    style = normalize_style(style)
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    style_config = PRESENTATION_STYLES[style]
    
    # Get persona-specific CTA content
    title_text = content.get("title", "Ready to Transform Your Business?")
    cta_text = content.get("cta_text", "Take the first step today")
    contact_info = content.get("contact_info", "Contact us to get started")
    
    # Line 1: Main title with adjusted position
    title_box = slide.shapes.add_textbox(
        Inches(1.0),
        Inches(2.0),  # Moved up for better spacing
        Inches(11.33),
        Inches(1.2)   # Reduced height
    )
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = truncate_text_for_slide(title_text, max_chars=60, max_words=12)  # Limit title length
    title_p.alignment = PP_ALIGN.CENTER
    
    for run in title_p.runs:
        run.font.name = style_config["fonts"]["heading"]
        run.font.size = Pt(style_config["font_sizes"]["title"])
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(style_config["colors"]["primary"])
    
    # Simple accent line between the texts with adjusted position
    accent_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2.0),
        Inches(3.5),  # Adjusted position
        Inches(9.33),
        Inches(0.05)
    )
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = RGBColor.from_string(style_config["colors"]["accent"])
    accent_line.line.fill.background()
    
    # Line 2: Call to action with adjusted position and size
    cta_box = slide.shapes.add_textbox(
        Inches(1.0),
        Inches(4.0),  # Adjusted position
        Inches(11.33),
        Inches(0.8)   # Reduced height
    )
    cta_tf = cta_box.text_frame
    cta_p = cta_tf.paragraphs[0]
    cta_p.text = truncate_text_for_slide(cta_text, max_chars=50, max_words=10)  # Limit CTA text length
    cta_p.alignment = PP_ALIGN.CENTER
    
    for run in cta_p.runs:
        run.font.name = style_config["fonts"]["body"]
        run.font.size = Pt(style_config["font_sizes"]["subtitle"])
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(style_config["colors"]["secondary"])
    
    # Line 3: Contact info with adjusted position
    contact_box = slide.shapes.add_textbox(
        Inches(1.0),
        Inches(5.0),  # Adjusted position
        Inches(11.33),
        Inches(0.6)   # Reduced height
    )
    contact_tf = contact_box.text_frame
    contact_p = contact_tf.paragraphs[0]
    contact_p.text = truncate_text_for_slide(contact_info, max_chars=40, max_words=8)  # Limit contact info length
    contact_p.alignment = PP_ALIGN.CENTER
    
    for run in contact_p.runs:
        run.font.name = style_config["fonts"]["body"]
        run.font.size = Pt(style_config["font_sizes"]["body"])
        run.font.italic = True
        run.font.color.rgb = RGBColor.from_string(style_config["colors"]["accent"])
    
    return slide

def fetch_image_with_cache(slide_type: str, context: Dict[str, Any], used_images: set, use_placeholders: bool = False):
    """Fetch an image with caching and deduplication to prevent repeating the same image."""
    import streamlit as st
    
    # Check if we're creating a new presentation or using a cached one
    is_new_presentation = context.get('is_new_presentation', True)
    
    # Check if we already have this image in cache
    if 'original_images_cache' in st.session_state and f'{slide_type}_slide' in st.session_state.original_images_cache:
        # Get the cached image data
        cached_data = st.session_state.original_images_cache[f'{slide_type}_slide']
        
        # Create a hash of the image data to track unique images
        img_hash = hashlib.md5(cached_data).hexdigest()
        
        # Check if this image has already been used in this presentation
        if img_hash not in used_images:
            # Mark image as used
            used_images.add(img_hash)
            # Return the cached image data
            return io.BytesIO(cached_data)
        elif not is_new_presentation:
            # If we're not creating a new presentation (using cache),
            # just return the cached image even if it's used before
            # This prevents unnecessary API calls
            return io.BytesIO(cached_data)
        else:
            # Image has been used before in a new presentation, fetch a different one
            return fetch_unique_image(slide_type, context, used_images, use_placeholders)
    
    # No cached image, fetch a new one only if creating a new presentation
    if is_new_presentation:
        return fetch_unique_image(slide_type, context, used_images, use_placeholders)
    else:
        # For cached presentations, don't fetch new images if not in cache
        return None

def fetch_unique_image(slide_type: str, context: Dict[str, Any], used_images: set, use_placeholders: bool = False):
    """Fetch a unique image that hasn't been used in the presentation yet."""
    # Don't fetch new images if we're not creating a new presentation
    if not context.get('is_new_presentation', True):
        return None
        
    # Try up to 3 times to get a unique image
    for _ in range(3):
        # Fetch a new image
        image_data = fetch_image_for_slide(slide_type, context, use_placeholders)
        if not image_data:
            return None
            
        # Check if unique
        image_data.seek(0)
        img_bytes = image_data.read()
        image_data.seek(0)
        
        import hashlib
        img_hash = hashlib.md5(img_bytes).hexdigest()
        
        if img_hash not in used_images:
            # Mark as used and cache
            used_images.add(img_hash)
            
            # Cache the image
            import streamlit as st
            if 'original_images_cache' in st.session_state:
                st.session_state.original_images_cache[f'{slide_type}_slide'] = img_bytes
            
            return image_data
    
    # If we couldn't get a unique image after 3 tries, use a generic one
    return None

def get_search_terms_for_slide(slide_type: str, context: Dict[str, Any]) -> str:
    """Generate specific search terms for each slide type based on content and persona.
    Returns a space-separated string of search terms optimized for image search."""
    search_terms = []
    
    # Get persona and style configuration
    persona = context.get('metadata', {}).get('persona', 'business') if context else 'business'
    persona = normalize_style(persona)
    style_config = PRESENTATION_STYLES[persona]
    
    # Map slide types to image categories
    slide_type_mapping = {
        'title_slide': 'title',
        'problem_slide': 'problem',
        'solution_slide': 'solution',
        'features_slide': 'features',
        'advantage_slide': 'advantage',
        'audience_slide': 'audience',
        'market_slide': 'market',
        'roadmap_slide': 'roadmap',
        'team_slide': 'team',
        'cta_slide': 'cta',
        'success_stories_slide': 'success',
        'impact_slide': 'impact',
        'future_vision_slide': 'future'
    }
    
    # Get the base slide type
    base_type = slide_type_mapping.get(slide_type, slide_type.replace('_slide', ''))
    
    # Get persona-specific image terms for this slide type
    if 'images' in style_config and base_type in style_config['images']:
        # Take only the first two image terms to keep search focused
        search_terms.extend(style_config['images'][base_type][:2])
    
    # Get product name and industry for additional context
    product_name = context.get("metadata", {}).get("product_name", "") if context else ""
    industry = context.get('metadata', {}).get('industry', "")
    
    # Add slide-specific terms if available in context
    if slide_type in context:
        slide_content = context[slide_type]
        
        # Extract keywords from title
        if 'title' in slide_content:
            title_words = [w for w in slide_content['title'].split() 
                         if len(w) > 3 and w.lower() not in ['the', 'and', 'for', 'with']]
            if title_words:
                search_terms.append(title_words[0])  # Add only the most relevant word
        
        # Extract keywords from description or paragraph
        if 'description' in slide_content or 'paragraph' in slide_content:
            text = slide_content.get('description', slide_content.get('paragraph', ''))
            words = [w for w in text.split() 
                    if len(w) > 4 and w.lower() not in ['their', 'there', 'these', 'those', 'about', 'because', 'through']]
            if words:
                search_terms.append(words[0])  # Add only the most relevant word
        
        # Extract keywords from bullets if available
        if 'bullets' in slide_content and isinstance(slide_content['bullets'], list):
            for bullet in slide_content['bullets'][:1]:  # Use only first bullet
                if isinstance(bullet, str):
                    words = [w for w in bullet.split() 
                            if len(w) > 4 and w.lower() not in ['their', 'there', 'these', 'those', 'about', 'because', 'through']]
                    if words:
                        search_terms.append(words[0])  # Add only the most relevant word
    
    # Clean up search terms
    cleaned_terms = []
    for term in search_terms:
        if isinstance(term, str):
            # Remove special characters and convert to lowercase
            cleaned = ' '.join(''.join(c for c in term.lower() if c.isalnum() or c.isspace()).split())
            if cleaned and len(cleaned) > 3:
                cleaned_terms.append(cleaned)
    
    # Add product name and industry if available
    if product_name:
        cleaned_product = ' '.join(''.join(c for c in product_name.lower() if c.isalnum() or c.isspace()).split())
        if cleaned_product:
            cleaned_terms.insert(0, cleaned_product)
    
    if industry:
        cleaned_industry = ' '.join(''.join(c for c in industry.lower() if c.isalnum() or c.isspace()).split())
        if cleaned_industry:
            cleaned_terms.append(cleaned_industry)
    
    # Remove duplicates while preserving order
    unique_terms = list(dict.fromkeys(cleaned_terms))
    
    # Limit to 3 terms and join with spaces
    final_terms = unique_terms[:3]
    
    # Add qualifiers for better image results
    qualifiers = {
        'title': 'professional business',
        'problem': 'business challenge',
        'solution': 'business solution',
        'features': 'modern business',
        'advantage': 'professional business',
        'audience': 'business people',
        'market': 'business market',
        'roadmap': 'business strategy',
        'team': 'professional team',
        'cta': 'business action',
        'success': 'business success',
        'impact': 'business impact',
        'future': 'business future'
    }
    
    # Add a relevant qualifier based on slide type
    if base_type in qualifiers:
        final_terms.insert(0, qualifiers[base_type])
    
    # Join terms with spaces and ensure proper formatting
    search_string = ' '.join(final_terms).strip()
    
    # Add general qualifiers for better image quality
    search_string += ' professional high quality business'
    
    return search_string

def create_market_slide_wrapper(prs: Presentation, content: Dict[str, Any], presentation_context: Dict[str, Any] = None, used_images: set = None):
    """
    Wrapper function for create_market_slide that ensures content fits within slide boundaries.
    """
    # Ensure market content is not too long
    if "description" in content:
        content["description"] = truncate_text_for_slide(content["description"], max_chars=200, max_words=35)
    
    try:
        from ppt_generator_additions import create_market_slide
        return create_market_slide(prs, content, presentation_context, used_images)
    except Exception as e:
        print(f"Error creating market slide: {e}")
        return create_fallback_slide(prs, "Market Analysis", content)

def create_roadmap_slide_wrapper(prs: Presentation, content: Dict[str, Any], presentation_context: Dict[str, Any] = None, used_images: set = None):
    """
    Wrapper function for create_roadmap_slide to ensure backward compatibility and error handling.
    """
    try:
        # Import the function from ppt_generator_additions
        from ppt_generator_additions import create_roadmap_slide
        return create_roadmap_slide(prs, content, presentation_context, used_images)
    except Exception as e:
        print(f"Error creating roadmap slide: {e}")
        # Fallback to a simple slide if there's an error
        return create_fallback_slide(prs, "Product Roadmap", content)

def create_team_slide_wrapper(prs: Presentation, content: Dict[str, Any], presentation_context: Dict[str, Any] = None, used_images: set = None):
    """
    Wrapper function for create_team_slide that only creates the slide when explicitly requested.
    """
    # Check if team slide is explicitly requested
    if not presentation_context.get("include_team_slide", False):
        return None
        
    # Check if team members are explicitly provided
    team_members = content.get("team_members", [])
    if not team_members:
        return None
        
    # Check if team members have real names (not generic placeholders)
    generic_names = {"John Doe", "Jane Smith", "Team Member", "CEO", "CTO", "Manager", 
                    "Chief Executive Officer", "Chief Technology Officer", "Chief Operating Officer",
                    "COO", "CFO", "CMO", "Director", "VP", "Vice President"}
    
    has_real_names = False
    for member in team_members:
        if isinstance(member, dict):
            name = member.get("name", "")
            title = member.get("title", "")
            # Check if both name and title are non-generic
            if name and name not in generic_names and title not in generic_names:
                has_real_names = True
                break
        elif isinstance(member, str) and member not in generic_names:
            has_real_names = True
            break
    
    if not has_real_names:
        return None
    
    try:
        from ppt_generator_additions import create_team_slide
        return create_team_slide(prs, content, presentation_context, used_images)
    except Exception as e:
        print(f"Error creating team slide: {e}")
        return None

def create_fallback_slide(prs: Presentation, title: str, content: Dict[str, Any]):
    """
    Create a simple fallback slide when specialized slide generation fails.
    
    Args:
        prs: Presentation object
        title: Title to display on the slide
        content: Original content dictionary for the slide
    """
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_left = Inches(0.5)
    title_top = Inches(0.5)
    title_width = Inches(12.33)
    title_height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.auto_size = False
    title_p = title_tf.paragraphs[0]
    title_p.text = content.get("title", title)
    title_p.alignment = PP_ALIGN.LEFT
    
    # Apply title formatting
    for run in title_p.runs:
        run.font.name = FONTS["title"]
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Content box
    left = Inches(0.5)
    top = Inches(1.7)
    width = Inches(12.33)
    height = Inches(5.0)
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.auto_size = False
    
    p = tf.paragraphs[0]
    
    # Try to extract some useful content from the content dict
    text = ""
    if "description" in content:
        text += content["description"]
    elif "paragraph" in content:
        text += content["paragraph"]
    elif "content" in content:
        text += content["content"]
    
    if not text and "bullets" in content and isinstance(content["bullets"], list):
        for bullet in content["bullets"][:5]:  # Limit to 5 items
            text += f"• {bullet}\n"
    
    if not text:
        text = "Content for this slide is being prepared."
    
    p.text = text
    
    # Format text
    for run in p.runs:
        run.font.name = FONTS["body"]
        run.font.size = Pt(20)
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    return slide

def create_bullet_list(slide, bullets, left, top, width, height, style="business", max_bullets=5, min_bullets=3):
    """Create a bullet list with standardized formatting."""
    style = normalize_style(style)
    style_config = PRESENTATION_STYLES[style]
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    
    for i, bullet_text in enumerate(bullets[:max_bullets]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        truncated_text = truncate_text_for_slide(bullet_text, max_chars=100, max_words=20)
        p.text = f"• {truncated_text}"
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(3)
        p.space_after = Pt(3)
        p.line_spacing = 1.0
        
        for run in p.runs:
            run.font.name = style_config["fonts"]["body"]
            run.font.size = CONTENT_DIMENSIONS['bullet_section']['font_size']
            run.font.color.rgb = RGBColor.from_string(style_config["colors"]["text"])
    
    return text_box

def determine_content_type(bullets):
    """Determine the type of content based on existing bullets."""
    combined_text = " ".join(str(b) for b in bullets).lower()
    
    if any(word in combined_text for word in ["market", "growth", "opportunity", "size"]):
        return "market"
    elif any(word in combined_text for word in ["solution", "feature", "platform", "product"]):
        return "solution"
    elif any(word in combined_text for word in ["challenge", "need", "require", "face"]):
        return "challenge"
    else:
        return "general"

def generate_additional_bullet(content_type):
    """Generate additional bullet points based on content type."""
    additional_bullets = {
        "market": [
            "Growing demand for integrated security solutions",
            "Increasing adoption of cloud-based security platforms",
            "Rising focus on regulatory compliance and data protection",
            "Expanding threat landscape driving security investments"
        ],
        "solution": [
            "Advanced threat detection and response capabilities",
            "Seamless integration with existing security tools",
            "Automated security workflows and orchestration",
            "Real-time monitoring and alerting system"
        ],
        "challenge": [
            "Need for comprehensive security visibility",
            "Resource constraints in security operations",
            "Complex compliance requirements management",
            "Rapid response to emerging threats"
        ],
        "general": [
            "Enhanced operational efficiency and productivity",
            "Improved security posture and risk management",
            "Streamlined compliance and reporting processes",
            "Cost-effective security management"
        ]
    }
    
    import random
    bullets = additional_bullets.get(content_type, additional_bullets["general"])
    return random.choice(bullets)

def expand_content_intelligently(text, content_type):
    """Intelligently expand content based on context and type."""
    if not text:
        return text
        
    # Define content-specific expansions
    expansions = {
        "problem": {
            "keywords": ["challenge", "issue", "risk", "threat", "difficulty"],
            "additions": [
                "This poses significant risks to business operations and data security.",
                "Organizations must address this challenge proactively to maintain compliance.",
                "The impact on overall security posture is substantial and growing.",
                "Traditional approaches have proven insufficient in today's threat landscape."
            ]
        },
        "solution": {
            "keywords": ["solution", "platform", "system", "approach"],
            "additions": [
                "Our solution delivers comprehensive protection across the entire attack surface.",
                "The platform integrates seamlessly with existing security infrastructure.",
                "Advanced analytics provide real-time insights and actionable intelligence.",
                "Automated responses minimize manual intervention and reduce response times."
            ]
        },
        "feature": {
            "keywords": ["feature", "capability", "function", "tool"],
            "additions": [
                "This capability enhances operational efficiency while reducing complexity.",
                "Users benefit from streamlined workflows and improved productivity.",
                "Integration capabilities maximize value across the security ecosystem.",
                "Regular updates ensure continued effectiveness against emerging threats."
            ]
        },
        "market": {
            "keywords": ["market", "industry", "sector", "opportunity"],
            "additions": [
                "Market trends indicate strong growth potential in the security sector.",
                "Industry adoption continues to accelerate as threats evolve.",
                "Regulatory requirements drive increased demand for comprehensive solutions.",
                "Digital transformation initiatives fuel expansion across all verticals."
            ]
        }
    }
    
    # Clean the text
    clean_text = text.strip()
    
    # Determine if we need to expand
    if len(clean_text) < 400:  # Allow for more substantial content
        # Find relevant expansions
        content_expansions = expansions.get(content_type, expansions["solution"])
        
        # Add relevant expansions
        for addition in content_expansions["additions"]:
            if len(clean_text) >= 400:
                break
            if not clean_text.endswith('.'):
                clean_text += "."
            clean_text += " " + addition
    
    return clean_text

def enhance_bullet_content(bullets, content_type):
    """Enhance bullet points with additional context and details."""
    enhanced = []
    min_bullet_length = 75  # Minimum length for each bullet
    
    for bullet in bullets:
        clean_text = bullet.strip()
        if len(clean_text) < min_bullet_length:
            # Add context-specific expansion
            expanded = expand_content_intelligently(clean_text, content_type)
            enhanced.append(expanded)
        else:
            enhanced.append(clean_text)
    
    return enhanced

def create_content_section(slide, content, left, top, width, height, style_config, content_type="general"):
    """Create a content section with standardized formatting."""
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    
    # Expand and truncate content appropriately
    expanded_content = expand_content_intelligently(content, content_type)
    truncated_content = truncate_text_for_slide(expanded_content, max_chars=250, max_words=50)
    
    p = tf.paragraphs[0]
    p.text = truncated_content
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(3)
    p.space_after = Pt(3)
    
    for run in p.runs:
        run.font.name = style_config["fonts"]["body"]
        run.font.size = CONTENT_DIMENSIONS['main_content']['font_size']
        run.font.color.rgb = RGBColor.from_string(style_config["colors"]["text"])
    
    return text_box

def add_slide_image(slide, image_data, content_start_y, style_config):
    """Add an image to a slide with proper positioning and fallback handling."""
    if image_data:
        pic = slide.shapes.add_picture(
            image_data,
            Inches(7.8),
            content_start_y + Inches(0.2),
            Inches(4.5),
            Inches(4.0)
        )
    else:
        # Add styled fallback icon
        icon_left = Inches(9.5)
        icon_top = content_start_y + Inches(1.0)
        icon_box = slide.shapes.add_textbox(icon_left, icon_top, Inches(2.0), Inches(2.0))
        icon_tf = icon_box.text_frame
        icon_tf.auto_size = False
        icon_p = icon_tf.paragraphs[0]
        icon_p.text = get_slide_icon("solution")["icon"]
        icon_p.alignment = PP_ALIGN.CENTER
        
        for run in icon_p.runs:
            run.font.name = style_config["fonts"]["body"]
            run.font.size = Pt(64)
            run.font.color.rgb = RGBColor.from_string(style_config["colors"]["primary"])

def create_presentation(content: Dict[str, Any], filename: str, image_manager=None, custom_slide_order: List[str] = None, style="business") -> io.BytesIO:
    """Create a PowerPoint presentation with persona-specific content and styling."""
    import streamlit as st
    import hashlib
    import json
    
    try:
        # Get the persona from content metadata or default to business
        persona = content.get('metadata', {}).get('persona', 'business')
        persona = normalize_style(persona)
        style = persona  # Use persona-specific style
        
        # Get the team configuration from session state
        team_config = st.session_state.get('team_config', {
            'include_team_slide': False,
            'team_members': []
        })
        
        # Create a cache key based on the content and configuration
        cache_key = None
        try:
            cache_input = {
                'content': content,
                'filename': filename,
                'custom_slide_order': custom_slide_order if custom_slide_order else [],
                'style': style,
                'persona': persona,
                'team_config': team_config
            }
            cache_key = hashlib.md5(json.dumps(cache_input, sort_keys=True).encode()).hexdigest()
            
            if 'ppt_cache' in st.session_state and cache_key in st.session_state.ppt_cache:
                cached_ppt = st.session_state.ppt_cache[cache_key]
                cached_ppt.seek(0)
                return cached_ppt
        except Exception as e:
            print(f"Cache key generation error: {e}")
        
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        style_config = PRESENTATION_STYLES[style]
        used_images = set()
        
        # Get persona-specific content
        persona_content = PERSONA_CONTENT.get(persona, PERSONA_CONTENT['business'])
        
        # Merge persona-specific content with user content
        for slide_type, slide_content in persona_content.items():
            if slide_type in content:
                try:
                    # Deep merge of content
                    merged_content = content[slide_type].copy()
                    for key, value in slide_content.items():
                        if key not in merged_content or not merged_content[key]:
                            merged_content[key] = value
                        elif isinstance(value, list) and isinstance(merged_content[key], list):
                            # Combine lists while removing duplicates
                            merged_content[key] = list(dict.fromkeys(merged_content[key] + value))
                    content[slide_type] = merged_content
                except Exception as e:
                    print(f"Error merging content for slide {slide_type}: {e}")
                    content[slide_type] = slide_content  # Use default content on error
        
        presentation_context = {
            'metadata': content.get('metadata', {}),
            'style': style,
            'style_config': style_config,
            'include_team_slide': team_config['include_team_slide'],
            'persona': persona
        }
        
        # Use custom_slide_order if provided, otherwise generate default order
        if custom_slide_order:
            print(f"Using custom slide order: {custom_slide_order}")
            slide_order = custom_slide_order
        else:
            # Default slide order logic
            core_slides = ['title_slide', 'problem_slide', 'solution_slide', 'features_slide']
            supporting_slides = ['advantage_slide', 'audience_slide']
            optional_slides = ['market_slide', 'roadmap_slide']
            
            if persona == 'technical':
                core_slides = ['title_slide', 'problem_slide', 'solution_slide', 'features_slide', 'advantage_slide']
                supporting_slides = ['roadmap_slide', 'audience_slide']
                optional_slides = ['market_slide']
            elif persona == 'executive':
                core_slides = ['title_slide', 'problem_slide', 'market_slide', 'solution_slide', 'advantage_slide']
                supporting_slides = ['features_slide', 'audience_slide']
                optional_slides = ['roadmap_slide']
            
            slide_order = []
            
            # Add slides in the default order
            for slide_type in core_slides + supporting_slides + optional_slides:
                if slide_type in content:
                    slide_order.append(slide_type)
            
            # Add team slide if configured
            if team_config['include_team_slide'] and team_config['team_members']:
                slide_order.append('team_slide')
            
            # Add CTA slide last if not in custom order
            if 'cta_slide' in content and 'cta_slide' not in slide_order:
                slide_order.append('cta_slide')
        
        print(f"Final slide order: {slide_order}")
        
        # Create slides in the specified order
        for slide_type in slide_order:
            if not slide_type in content:
                print(f"Warning: Slide {slide_type} not found in content, skipping...")
                continue
                
            slide_data = content[slide_type]
            if slide_type == 'title_slide':
                create_title_slide(prs, slide_data, style)
            elif slide_type == 'problem_slide':
                create_problem_slide(prs, slide_data, presentation_context, used_images)
            elif slide_type == 'solution_slide':
                create_solution_slide(prs, slide_data, presentation_context, used_images)
            elif slide_type == 'features_slide':
                create_features_slide(prs, slide_data, style)
            elif slide_type == 'advantage_slide':
                create_advantage_slide(prs, slide_data, presentation_context, used_images)
            elif slide_type == 'audience_slide':
                create_audience_slide(prs, slide_data, presentation_context, used_images)
            elif slide_type == 'market_slide':
                create_market_slide_wrapper(prs, slide_data, presentation_context, used_images)
            elif slide_type == 'roadmap_slide':
                create_roadmap_slide_wrapper(prs, slide_data, presentation_context, used_images)
            elif slide_type == 'team_slide':
                create_team_slide_wrapper(prs, slide_data, presentation_context, used_images)
            elif slide_type == 'cta_slide':
                create_cta_slide(prs, slide_data, style)
        
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        
        try:
            if cache_key:
                if 'ppt_cache' not in st.session_state:
                    st.session_state.ppt_cache = {}
                cache_copy = io.BytesIO(output.getvalue())
                st.session_state.ppt_cache[cache_key] = cache_copy
        except Exception as e:
            print(f"Error caching presentation: {e}")
        
        return output
    except Exception as e:
        print(f"Error creating presentation: {e}")
        raise e

def truncate_text_for_slide(text, max_chars=None, max_words=None):
    """Intelligently truncate text for slides while maintaining complete sentences."""
    if not text:
        return ""
    
    # Set professional default limits for slides
    max_chars = max_chars or 100  # Allow more characters for proper content
    max_words = max_words or 20   # Allow more words for complete thoughts
    
    # Clean the text first
    text = text.strip()
    
    # Check if text needs truncation
    words = text.split()
    if len(words) <= max_words and len(text) <= max_chars:
        return ensure_complete_sentences(text)
    
    # Truncate by words first
    if len(words) > max_words:
        text = " ".join(words[:max_words])
    
    # Then check character length
    if len(text) > max_chars:
        # Find the last complete sentence within the limit
        last_period = text.rfind('.', 0, max_chars)
        last_exclamation = text.rfind('!', 0, max_chars)
        last_question = text.rfind('?', 0, max_chars)
        
        # Get the last sentence end position
        last_sentence_end = max(last_period, last_exclamation, last_question)
        
        if last_sentence_end > 0:  # If we found a sentence end
            text = text[:last_sentence_end + 1].strip()
        else:
            # Find the last complete word
            last_space = text.rfind(' ', 0, max_chars)
            if last_space > 0:
                text = text[:last_space].strip()
            else:
                text = text[:max_chars].strip()
    
    # Ensure proper ending
    return ensure_complete_sentences(text)

def ensure_complete_sentences(text):
    """Ensure text ends with complete sentences and proper punctuation."""
    if not text:
        return text
        
    # Remove trailing spaces and ellipsis
    text = text.rstrip('. ').rstrip('...')
    
    # If text ends with alphanumeric character, add a period
    if text[-1].isalnum():
        text += "."
    
    # Common sentence endings to check
    endings = {
        'manua': 'manual processes.',
        'vulnerab': 'vulnerability management.',
        'efficien': 'efficiency.',
        'secur': 'security measures.',
        'complian': 'compliance requirements.',
        'detect': 'detection capabilities.',
        'automa': 'automation capabilities.',
        'respon': 'response procedures.',
        'sophisti': 'sophisticated approach.',
        'integra': 'integration capabilities.',
        'platfor': 'platform features.',
        'soluti': 'solution benefits.',
        'analy': 'analysis capabilities.',
        'monitor': 'monitoring system.',
        'protect': 'protection measures.'
    }
    
    # Check for incomplete words and complete them
    last_word = text.split()[-1].lower()
    for partial, complete in endings.items():
        if last_word.startswith(partial) and last_word != complete.split()[0].lower():
            text = " ".join(text.split()[:-1]) + " " + complete
            break
    
    return text

def add_fallback_icon(slide, style_config, slide_type, icon_left, icon_top):
    """
    Add a fallback icon when an image is not available or fails to load.
    
    Args:
        slide: The PowerPoint slide object
        style_config: The style configuration dictionary
        slide_type: The type of slide ('impact', 'success_stories', etc.)
        icon_left: Left position of the icon (in inches)
        icon_top: Top position of the icon (in inches)
    """
    # Create text box for icon
    icon_box = slide.shapes.add_textbox(icon_left, icon_top, Inches(2.0), Inches(2.0))
    icon_tf = icon_box.text_frame
    icon_tf.auto_size = False
    icon_p = icon_tf.paragraphs[0]
    
    # Choose appropriate icon based on slide type
    icons = {
        'impact': "📈",      # Chart increasing
        'success_stories': "🏆",  # Trophy
        'future_vision': "🚀",    # Rocket
        'market': "📊",      # Bar chart
        'features': "⚡",     # Lightning bolt
        'solution': "🎯",     # Target
        'advantage': "💪",    # Strong arm
        'audience': "👥",     # People
        'problem': "❗",      # Exclamation
        'default': "📋"      # Clipboard
    }
    
    # Get icon or use default if not found
    icon_p.text = icons.get(slide_type, icons['default'])
    
    # Center align the icon
    icon_p.alignment = PP_ALIGN.CENTER
    
    # Apply icon formatting
    for run in icon_p.runs:
        run.font.name = style_config["fonts"]["body"]
        run.font.size = Pt(64)
        run.font.color.rgb = RGBColor.from_string(style_config["colors"]["primary"])
    
    # Add subtle accent line below icon
    line_left = icon_left - Inches(0.5)
    line_top = icon_top + Inches(2.2)
    line_width = Inches(3.0)
    line_height = Inches(0.02)
    
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        line_left,
        line_top,
        line_width,
        line_height
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor.from_string(style_config["colors"]["accent"])
    line.line.fill.background()

def create_alternative_slide(prs: Presentation, content: Dict[str, Any], slide_type: str, presentation_context: Dict[str, Any]):
    """Create alternative slides with proper content and consistent formatting."""
    style = normalize_style(presentation_context.get('style', 'business'))
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    style_config = PRESENTATION_STYLES[style]
    
    # Get content from template
    selected_content = SLIDE_CONTENT_TEMPLATES.get(slide_type, SLIDE_CONTENT_TEMPLATES['impact']).copy()
    
    # If custom content is provided, merge it with template
    if isinstance(content, dict):
        if 'title' in content:
            selected_content['title'] = content['title']
        if 'bullets' in content and isinstance(content['bullets'], list):
            # For last slide, keep it focused but substantial
            bullets = content['bullets'][:3]  # Take up to 3 key points
            selected_content['bullets'] = [truncate_text_for_slide(b, max_chars=100, max_words=20) for b in bullets]
    
    # Create section header
    content_start_y = create_section_header(slide, selected_content['title'], style)
    
    # Create bullet list with proper content
    create_bullet_list(
        slide,
        selected_content['bullets'],
        left=Inches(0.5),
        top=content_start_y + Inches(0.3),
        width=Inches(7.0),
        height=Inches(4.5),  # Increased height for content
        style=style,
        max_bullets=4,  # Allow up to 4 bullets for main slides
        min_bullets=3   # Minimum 3 bullets for substance
    )
    
    # Add image
    image_left = Inches(8.0)
    image_top = content_start_y + Inches(0.3)
    image_width = Inches(4.5)
    image_height = Inches(4.5)  # Match content height
    
    # Try to get image from content
    image_data = None
    if isinstance(content, dict) and 'image' in content:
        try:
            image_data = content['image']
        except Exception:
            pass
    
    if image_data:
        try:
            pic = slide.shapes.add_picture(
                image_data,
                image_left,
                image_top,
                image_width,
                image_height
            )
        except Exception:
            add_fallback_icon(slide, style_config, slide_type, icon_left=Inches(9.25), 
                            icon_top=content_start_y + Inches(1.0))
    else:
        add_fallback_icon(slide, style_config, slide_type, icon_left=Inches(9.25), 
                         icon_top=content_start_y + Inches(1.0))
    
    return slide

def generate_alternative_content(slide_type: str) -> Dict[str, Any]:
    """Generate alternative content for different slide types."""
    return SLIDE_CONTENT_TEMPLATES.get(slide_type, SLIDE_CONTENT_TEMPLATES['impact'])